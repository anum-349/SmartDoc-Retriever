import os
import json
import hashlib
from typing import List

import google.generativeai as genai
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader, CSVLoader, Docx2txtLoader
from langchain_core.documents import Document
from pptx import Presentation
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_community.llms import GooglePalm
from langchain.chains import RetrievalQA
from langchain_google_genai import ChatGoogleGenerativeAI
import pdfplumber
from dotenv import load_dotenv

load_dotenv()

# Load API key from .env file
os.environ['GOOGLE_API_KEY'] = os.getenv("GOOGLE_API_KEY")
try:
    genai.configure(api_key=os.environ['GOOGLE_API_KEY'])
except Exception as e:
    print("Error Configuring API: ", e)

llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash")

INDEX_METADATA_FILE = "faiss_metadata.json"
INDEX_FOLDER = "faiss_index"

def compute_file_hash(file_path, chunk_size=8192):
    sha256 = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            while True:
                 # Standard read operation
                chunk = f.read(chunk_size) 
                if not chunk:
                    break # Exit the loop when done reading
                sha256.update(chunk)
        return sha256.hexdigest()
    except Exception as e:
        print(f"❌ Error computing hash for {file_path}: {e}")
        return ""

def load_index_and_metadata(embedding_model):
    vector_store = None
    if os.path.exists(INDEX_FOLDER):
        print(f"Loading existing FAISS index from {INDEX_FOLDER}...")
        try:
            vector_store = FAISS.load_local(INDEX_FOLDER, embedding_model, allow_dangerous_deseialization=True)        
        except Exception as e:
            print(f"❌ Error loading FAISS index: {e}. Rebuilding index.")

    metadata = {}
    if os.path.exists(INDEX_METADATA_FILE):
        try:
            with open(INDEX_METADATA_FILE, "r") as f:
                metadata = json.load(f)
        except Exception as e:
            print(f"❌ Error loading metadata file: {e}. Starting fresh metadata.")
        
    return vector_store, metadata

def save_index_and_metadata(vector_store, metadata):
    if vector_store:
        print(f"Saving FAISS index to {INDEX_FOLDER}...")
        vector_store.save_local(INDEX_FOLDER)
    
    print(f"Saving Metadata to {INDEX_METADATA_FILE}...")
    with open(INDEX_METADATA_FILE, "w") as f:
        json.dump(metadata, f, indent=2)

def update_indexx_with_folder(folder_path, embedding_model, text_splitter):
    vector_store, metadata = load_index_and_metadata(embedding_model)
    new_docs = []

    current_files = {}
    for root, _, files in os.walk(folder_path):
        for  file in files:
            file_path = os.path.join(root, file)
            if file.startswith('.') or file.startswith('~'):
                continue
            file_hash = compute_file_hash(file_path)
            last_modified = os.path.getmtime(file_path)
            current_files[file_path] = {"hash": file_hash, "last_modified": last_modified}

    deleted_files = [f for f in metadata if f not in current_files]
    if deleted_files:
        print(f"Removing {len(deleted_files)} delete file from index")
        vector_store = None
        metadata = {}

    files_to_index = []
    for file_path, file_info in current_files.items():
        if file_path in  metadata and metadata[file_path]["hash"] == file_info["hash"]:
            continue

        print(f"Process {os.path.basename(file_path)}")
        files_to_index.append(file_path)

    if vector_store is None or files_to_index:
        if vector_store is None:
            print("Performing full index rebuild...")
            all_files = list(current_files.keys())
        else:
            all_files = files_to_index

        all_documents = []

        for file_path in all_files:
            docs = load_single_document(file_path)
            if(docs):
                chunks = text_splitter.split_documents(docs)
                if chunks:
                    all_documents.extend(chunks)

                    if file_path in current_files:
                        metadata[file_path] = current_files[file_path]

        if not all_documents:
            print("No documents loaded or chunks generated. Index remains unchanged.")
            return vector_store, []
        
        if vector_store is None:
            print(f"Building new FAISS index with {len(all_documents)} chunks.")
            vector_store = FAISS.from_documents(all_documents, embedding_model)
        else:
            print(f"Adding {len(all_documents)} new chunks to existing FAISS index.")
            vector_store.add_documents(all_documents)

        save_index_and_metadata(vector_store, metadata)
        new_chunks = all_documents

    else:
        print("Index is up to date. No files modified or deleted.")
    
    return vector_store, new_chunks

def load_pptx(file_path):
    """Load a PowerPoint file and return Document objects."""
    documents = []
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        content = "\n".join(text)
        documents.append(Document(page_content=content, metadata={"source": file_path}))
    except Exception as e:
        print(f"❌ Error loading PowerPoint file {file_path}: {e}")
    return documents

def load_single_document(file_path):
    documents = []
    file = os.path.basename(file_path)
    try:
        if file.endswith(".pdf"):
            try:
                with pdfplumber.open(file_path) as pdf:
                    docs = [Document(page_content=page.extract_text(), 
                                    metadata={"source": file_path, "page": i + 1}) 
                for i, page in enumerate(pdf.pages) if page.extract_text()]
                documents.extend([doc for doc in docs if doc.page_content.strip()])
            except Exception as e:
                print(f"❌ Error reading PDF {file}: {e}")
                docs = []  # Return an empty list in case of an error                    
        elif file.endswith(".txt"):
            documents = TextLoader(file_path, encoding='utf-8') .load()               
        elif file.endswith(".docx"):
            documents = Docx2txtLoader(file_path).load()
        elif file.endswith(".csv"):
            documents = CSVLoader(file_path).load()
        elif file.endswith(".pptx"):
            documents = load_pptx(file_path)
        else:
            print(f"⚠️ Unsupported file: {file}")
    except Exception as e:
        print(f"❌ Error loading {file}: {str(e)}")

    return documents

folder_path = "./Software-engineering"  # Change this if needed

if not os.path.exists(folder_path):
    print(f"❌ Error: The folder '{folder_path}' does not exist.")
    exit()

if len(os.listdir(folder_path)) == 0:
    print("❌ Error: No files found in the folder.")
    exit()

print("load/process documents and update index...")
    
text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)

# Convert text to embeddings
print('Initializing HuggingFace Embeddings...')
embedding_model = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2",
    cache_folder="C:/Users/akraj/.cache/huggingface/hub",
    model_kwargs={}  # Keep this empty if not needed
)

vector_store, chunks = update_indexx_with_folder(folder_path, embedding_model, text_splitter)

if vector_store is None:
    print("Error: Vector store could not be created or loaded.")
    exit()
    
# Function to answer queries using RAG
def answer_query(query):
    if vector_store is None:
        return "Error: Document index is not available."
    
    retriever = vector_store.as_retriever(search_kwargs={"k": 5})
    qa_chain = RetrievalQA.from_chain_type(llm, retriever=retriever)
    try:
        response = qa_chain.invoke({"query": query})
        return response.get("result", "No answer found.")
    except Exception as e:
        return f"An error occurred during query: {e}"

# Example user interaction
if __name__ == "__main__":
    while True:
        user_query = input("Ask a question (or type 'exit' to quit): ")
        if user_query.lower() == 'exit':
            break
        print("Thinking...")
        response = answer_query(user_query)
        print("\nChatbot: ", response, "\n")
